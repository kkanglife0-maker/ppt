import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from utils import load_json

class SlideFactory:
    def __init__(self, prs, style_config, hospital_data):
        self.prs = prs
        self.style = style_config
        self.data = hospital_data
        self.points_to_inch = 1/72
        self.px_to_inch = 1/96 # 96 DPI standard
        
    def _normalize_color(self, hex_color):
        """#FFF -> FFFFFF, #666 -> 666666, #123456 -> 123456"""
        hex_color = hex_color.replace("#", "")
        if len(hex_color) == 3:
            return "".join([c*2 for c in hex_color])
        return hex_color

    def _add_top_accent_bar(self, slide):
        """슬라이드 상단에 포인트 블루 액센트 바 추가"""
        bar_height = 12 # px
        left = 0
        top = 0
        width = self.prs.slide_width
        height = Inches(bar_height * self.px_to_inch)
        
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(self._normalize_color(self.style['colors']['point_blue']))
        shape.line.fill.background()

    def _add_footer_and_slide_number(self, slide, slide_num):
        """슬라이드 하단 정보 및 페이지 번호 추가"""
        # 푸터 텍스트
        footer_text = "아벨 전략팀 | 전략기획실"
        left = Inches(1)
        top = Inches(7)
        width = Inches(5)
        height = Inches(0.3)
        
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        p = tx_box.text_frame.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor.from_string("888888")
        
        # 페이지 번호
        page_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
        p = page_box.text_frame.paragraphs[0]
        p.text = str(slide_num)
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT
        
    def add_title_slide(self, slide_data):
        """1장 표지 생성"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank layout
        
        # 배경색 설정 (도형으로 채우기)
        # self._set_background(slide)
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "AEO 제안서")
        p.font.size = Pt(self.style['font_sizes']['title'])
        p.font.bold = True
        p.font.name = self.style['fonts']['header']
        p.font.color.rgb = RGBColor.from_string(self.style['colors']['point_blue'])
        
        # 부제목
        subtitle_text = slide_data.get("subtitle", "").format(hospital_name=self.data.get("hospital_name", "본 병원"))
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(self.style['font_sizes']['subtitle'])
        p.font.name = self.style['fonts']['body']
        
        # 하단 텍스트 (전략기획실)
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.33), Inches(0.5))
        p = footer_box.text_frame.paragraphs[0]
        p.text = slide_data.get("footer", "아벨 AI 마케팅 [전략기획실]")
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT

    def add_bullet_slide(self, slide_data):
        """고정형 불렛 포인트 슬라이드"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1]) # Title and Content
        slide.shapes.title.text = slide_data.get("title", "")
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "" # Clear default
        
        for item in slide_data.get("content", []):
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    def add_object_slide(self, slide_data):
        """객체 기반 디자인 슬라이드 생성 (HTML 변환용)"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank layout
        
        # 96 DPI conversion (1280px = 13.333 inches)
        px_to_inch = 1/96
        
        for obj in slide_data.get("objects", []):
            obj_type = obj.get("type")
            left = Inches(obj.get("left", 0) * px_to_inch)
            top = Inches(obj.get("top", 0) * px_to_inch)
            width = Inches(obj.get("width", 0) * px_to_inch)
            height = Inches(obj.get("height", 0) * px_to_inch)
            
            if obj_type == "shape":
                shape_type = MSO_SHAPE.RECTANGLE
                if obj.get("border_radius"):
                    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                
                shape = slide.shapes.add_shape(shape_type, left, top, width, height)
                
                # 배경색
                bg_color = obj.get("background_color", "#FFFFFF")
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(bg_color.replace("#", ""))
                
                # 테두리 설정
                if obj.get("border") == "none":
                    shape.line.fill.background()
                elif "border_color" in obj:
                    shape.line.color.rgb = RGBColor.from_string(obj["border_color"].replace("#", ""))
                else:
                    shape.line.fill.background()

            elif obj_type == "textbox":
                tx_box = slide.shapes.add_textbox(left, top, width, height)
                tf = tx_box.text_frame
                tf.word_wrap = True
                
                paragraphs = obj.get("paragraphs", [])
                if not paragraphs and "text" in obj:
                    paragraphs = [{"text": obj["text"]}]
                
                for i, p_data in enumerate(paragraphs):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
    def _get_nested_value(self, data, key_path):
        """'ai_tests[0].gpt_result' 식의 문자열 경로를 따라 데이터를 추출합니다."""
        import re
        parts = re.split(r'\.|\[|\]', key_path)
        parts = [p for p in parts if p]
        
        curr = data
        try:
            for p in parts:
                if p.isdigit():
                    curr = curr[int(p)]
                else:
                    curr = curr[p]
            return curr
        except (KeyError, IndexError, TypeError):
            return f"{{{{Missing: {key_path}}}}}"

    def _format_text_with_data(self, text):
        """텍스트 내의 {key.path} 또는 {key[0].path} 패턴을 병원 데이터로 치환합니다."""
        import re
        # {key...} 패턴 추출
        pattern = r'\{([^}]+)\}'
        matches = re.findall(pattern, text)
        
        for m in matches:
            val = self._get_nested_value(self.data, m)
            text = text.replace(f"{{{m}}}", str(val))
        return text

    def _add_background_image(self, slide, slide_num):
        """Render된 배경 이미지를 슬라이드 배경으로 삽입"""
        hospital_id = self.data.get("hospital_id")
        bg_path = os.path.join("1_assets", "hospitals", hospital_id, "renders", f"bg_{slide_num:02d}.png")
        
        if os.path.exists(bg_path):
            slide.shapes.add_picture(bg_path, 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
        else:
            print(f"Warning: Background image not found: {bg_path}")
            self._add_top_accent_bar(slide)

    def add_object_slide(self, slide_data):
        """객체 기반 디자인 슬라이드 생성 (배경 이미지 + 텍스트박스)"""
        slide_num = slide_data.get("slide_number", 0)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank layout
        
        # 1. 배경 이미지 삽입 (HTML 렌더링 결과)
        self._add_background_image(slide, slide_num)
        
        # 2. 텍스트 및 편집 가능한 객체만 추가
        px_to_inch = 1/96
        
        for obj in slide_data.get("objects", []):
            obj_type = obj.get("type")
            left = Inches(obj.get("left", 0) * px_to_inch)
            top = Inches(obj.get("top", 0) * px_to_inch)
            width = Inches(obj.get("width", 0) * px_to_inch)
            height = Inches(obj.get("height", 0) * px_to_inch)
            
            if obj_type == "shape":
                shape_type = MSO_SHAPE.RECTANGLE
                if obj.get("border_radius"):
                    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                
                shape = slide.shapes.add_shape(shape_type, left, top, width, height)
                
                # 배경색 및 투명도
                bg_color = obj.get("background_color", "#FFFFFF")
                opacity = obj.get("opacity", 1.0)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(self._normalize_color(bg_color))
                if opacity < 1.0:
                    shape.fill.transparency = 1.0 - opacity
                
                # 테두리 설정
                if obj.get("border") == "none":
                    shape.line.fill.background()
                elif "border_color" in obj:
                    shape.line.color.rgb = RGBColor.from_string(self._normalize_color(obj["border_color"]))
                    shape.line.width = Pt(obj.get("border_width", 1))
                else:
                    shape.line.fill.background()
                
                # 그림자 (Shadow) - 프리미엄 효과
                if obj.get("shadow"):
                    shadow = shape.shadow
                    shadow.inherit = False
                    shadow.visible = True
                    # PPT의 기본 바깥쪽 그림자 스타일 모사
                    shadow.blur_radius = Pt(obj.get("shadow_blur", 10))
                    shadow.distance = Pt(obj.get("shadow_distance", 5))
                    shadow.transparency = 1.0 - obj.get("shadow_opacity", 0.5)
                    shadow.color.rgb = RGBColor.from_string("000000")

            elif obj_type == "textbox":
                tx_box = slide.shapes.add_textbox(left, top, width, height)
                tf = tx_box.text_frame
                tf.word_wrap = True
                
                paragraphs = obj.get("paragraphs", [])
                if not paragraphs and "text" in obj:
                    paragraphs = [{"text": obj["text"]}]
                
                for i, p_data in enumerate(paragraphs):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    raw_text = p_data.get("text", "")
                    p.text = self._format_text_with_data(raw_text)
                    
                    # 폰트 설정
                    p.font.name = "Noto Sans KR"
                    p.font.size = Pt(p_data.get("font_size", 14))
                    
                    if p_data.get("font_weight", 400) >= 700:
                        p.font.bold = True
                    
                    color_hex = p_data.get("color", "#111111")
                    p.font.color.rgb = RGBColor.from_string(self._normalize_color(color_hex))
                    
                    if "alignment" in p_data:
                        align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "left": PP_ALIGN.LEFT}
                        p.alignment = align_map.get(p_data["alignment"], PP_ALIGN.LEFT)
                    
                    # 줄 간격 (Line Spacing)
                    # p.space_before = Pt(p_data.get("margin_top", 0))
                    # p.space_after = Pt(p_data.get("margin_bottom", 0))
                    
                    # 텍스트 내 특정 부분 하이라이트 (runs)
                    if "highlights" in p_data:
                        full_text = p.text
                        highlights = p_data["highlights"]
                        p.text = "" 
                        current_pos = 0
                        # 텍스트 데이터 보간 후의 위치를 찾아야 함
                        interpolated_full_text = self._format_text_with_data(raw_text)
                        
                        for h in highlights:
                            target_text = self._format_text_with_data(h["text"])
                            start_idx = interpolated_full_text.find(target_text, current_pos)
                            if start_idx == -1: continue
                            
                            if start_idx > current_pos:
                                run = p.add_run()
                                run.text = interpolated_full_text[current_pos:start_idx]
                                run.font.name = "Noto Sans KR"
                            
                            run = p.add_run()
                            run.text = target_text
                            run.font.name = "Noto Sans KR"
                            run.font.bold = h.get("font_weight", 400) >= 700
                            if "color" in h:
                                run.font.color.rgb = RGBColor.from_string(self._normalize_color(h["color"]))
                            
                            current_pos = start_idx + len(target_text)
                        
                        if current_pos < len(interpolated_full_text):
                            run = p.add_run()
                            run.text = interpolated_full_text[current_pos:]
                            run.font.name = "Noto Sans KR"

            elif obj_type == "icon":
                tx_box = slide.shapes.add_textbox(left, top, width, height)
                p = tx_box.text_frame.paragraphs[0]
                
                icon_map = {
                    "fa-keyboard": "⌨", "fa-list-ul": "•", "fa-hand-pointer": "👆",
                    "fa-arrow-down": "↓", "fa-wand-magic-sparkles": "✨",
                    "fa-circle-question": "❓", "fa-microchip": "💻",
                    "fa-clipboard-check": "✅", "fa-clock-rotate-left": "🕒",
                    "fa-arrow-right": "→", "fa-chart-pie": "📊",
                    "fa-arrows-rotate": "🔄", "fa-google": "G", "fa-robot": "🤖",
                    "fa-circle-user": "👤", "fa-angles-down": "︾",
                    "fa-clipboard-question": "📋", "fa-bolt": "⚡", "fa-code": "💻",
                    "fa-quote-left": "\"", "fa-quote-right": "\"",
                    "fa-magnifying-glass": "🔍", "fa-computer-mouse": "🖱",
                    "fa-globe": "🌐", "fa-comments": "💬", "fa-list-check": "✅"
                }
                fa_class = obj.get("icon_class", "").split(" ")[-1] # last class usually has icon name
                p.text = icon_map.get(fa_class, "•")
                p.font.size = Pt(obj.get("font_size", 24) * 0.75)
                color_hex = obj.get("color", "#666666").replace("#", "")
                p.font.color.rgb = RGBColor.from_string(color_hex)
                p.alignment = PP_ALIGN.CENTER

            elif obj_type == "chart":
                chart_data = CategoryChartData()
                data_list = obj.get("data", [])
                chart_data.categories = [d["name"] for d in data_list]
                chart_data.add_series("Series 1", [d["value"] for d in data_list])
                
                chart_type = XL_CHART_TYPE.PIE
                if obj.get("chart_type") == "doughnut":
                    chart_type = XL_CHART_TYPE.DOUGHNUT
                
                chart_shape = slide.shapes.add_chart(
                    chart_type, left, top, width, height, chart_data
                )
                chart = chart_shape.chart
                
                # 범례 설정
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                
                # 데이터 레이블
                chart.plots[0].has_data_labels = True
                data_labels = chart.plots[0].data_labels
                data_labels.show_percentage = True
                data_labels.show_category_name = True
                
                # 색상 설정 (Optional: python-pptx에서 차트 개별 색상 설정은 복잡하므로 기본 테마 사용 권장)

    def add_dynamic_slide(self, slide_num):
        """동적 슬라이드 생성 로직 분기 - 템플릿 우선 적용"""
        dynamic_template_path = os.path.join("3_templates", "dynamic_layouts", f"slide_{slide_num:02d}.json")
        
        # 04, 05, 06번은 각자 고유 템플릿 사용 (데이터 인덱스 때문)
        if os.path.exists(dynamic_template_path):
            slide_data = load_json(dynamic_template_path)
            print(f"Using dynamic template for slide {slide_num}")
            self.add_object_slide(slide_data)
            
            # 이미지 추가 (objects 레이아웃 뒤에 얹기)
            slide = self.prs.slides[-1]
            if slide_num in [4, 5, 6, 8, 9]:
                # 8, 9번도 이미지 있을 수 있으므로 처리
                self._add_image_if_exists(slide, slide_num, left=1, top=3, width=7)
            
            self._add_footer_and_slide_number(slide, slide_num)
            return

        method_name = f"create_slide_{slide_num:02d}"
        if hasattr(self, method_name):
            getattr(self, method_name)()
        else:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
            self._add_top_accent_bar(slide)
            self._add_footer_and_slide_number(slide, slide_num)
            slide.shapes.title.text = f"Slide {slide_num} (To be implemented)"
            
    def _add_image_if_exists(self, slide, slide_num, left=5, top=2, width=3):
        """이미지가 설정되어 있으면 삽입"""
        image_mapping = self.data.get("slide_images", {})
        image_names = image_mapping.get(str(slide_num), [])
        
        if not image_names:
            return
            
        hospital_id = self.data.get("hospital_id")
        for i, img_name in enumerate(image_names):
            img_path = os.path.join("1_assets", "hospitals", hospital_id, img_name)
            if os.path.exists(img_path):
                # 가로로 나열
                slide.shapes.add_picture(img_path, Inches(left + i*3.2), Inches(top), width=Inches(width))
            else:
                print(f"Warning: Image not found: {img_path}")
                # 빈칸 또는 안내 상자로 대체 (사용자 요청: 빈칸)
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + i*3.2), Inches(top), Inches(width), Inches(width * 0.56))
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string("F9F9F9")
                shape.line.color.rgb = RGBColor.from_string("DDDDDD")
                
                tx = shape.text_frame
                p = tx.paragraphs[0]
                p.text = "이미지 준비 중"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor.from_string("CCCCCC")
                p.alignment = PP_ALIGN.CENTER

    # --- 동적 슬라이드 구체적 구현 ---
    
    def create_slide_02(self):
        """2장 Executive Summary"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 2)
        slide.shapes.title.text = "Executive Summary"
        tf = slide.placeholders[1].text_frame
        
        summary = self.data.get("ai_analysis_summary", {}).get("current_status", "")
        p = tf.add_paragraph()
        p.text = f"현황: {summary}"
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "주요 현안:"
        for issue in self.data.get("top_issues", [])[:3]:
            sp = tf.add_paragraph()
            sp.text = f"• {issue}"
            sp.level = 1
            
        p = tf.add_paragraph()
        p.text = "실행 전략 요약:"
        for plan in self.data.get("implementation_plan", {}).get("summary", [])[:2]:
            sp = tf.add_paragraph()
            sp.text = f"• {plan}"
            sp.level = 1

    def create_slide_03(self):
        """3장 현재 AI 인식 상태 요약"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 3)
        slide.shapes.title.text = "현재 AI 인식 상태 분석"
        self._add_image_if_exists(slide, 3)
        
        tf = slide.placeholders[1].text_frame
        tf.text = "테스트 결과 종합:"
        for test in self.data.get("ai_tests", []):
            p = tf.add_paragraph()
            p.text = f"- {test['question']}: GPT({test['gpt_result']}), Gemini({test['gemini_result']}), Perplexity({test['perplexity_result']})"
            p.level = 1

    def create_slide_04(self): self._create_test_slide(4, 0)
    def create_slide_05(self): self._create_test_slide(5, 1)
    def create_slide_06(self): self._create_test_slide(6, 2)

    def _create_test_slide(self, slide_num, test_idx):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, slide_num)
        tests = self.data.get("ai_tests", [])
        if test_idx < len(tests):
            test = tests[test_idx]
            slide.shapes.title.text = f"AI 질문 테스트 {test_idx + 1}: {test['question_type']}"
            tf = slide.placeholders[1].text_frame
            tf.text = f"질문: {test['question']}"
            p = tf.add_paragraph()
            p.text = f"분석: {test['analysis']}"
            
            self._add_image_if_exists(slide, slide_num, top=4)
        else:
            slide.shapes.title.text = f"AI 질문 테스트 {test_idx + 1}"
            slide.placeholders[1].text_frame.text = "데이터 없음"

    def create_slide_07(self):
        """7장 왜 나오고 왜 안 나오는가관련 분석"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 7)
        slide.shapes.title.text = "AI 노출 로직 심층 분석"
        tf = slide.placeholders[1].text_frame
        
        summary = self.data.get("ai_analysis_summary", {})
        p = tf.add_paragraph()
        p.text = "노출 요인 (Strength):"
        for item in summary.get("reason_visible", []):
            sp = tf.add_paragraph(); sp.text = f"• {item}"; sp.level = 1
            
        p = tf.add_paragraph()
        p.text = "미노출 요인 (Weakness):"
        for item in summary.get("reason_not_visible", []):
            sp = tf.add_paragraph(); sp.text = f"• {item}"; sp.level = 1

    def create_slide_08(self):
        """8장 경쟁 병원 비교"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 8)
        slide.shapes.title.text = "경쟁 병원 AEO 현황 비교"
        self._add_image_if_exists(slide, 8, top=4)
        
        tf = slide.placeholders[1].text_frame
        tf.text = ""
        for comp in self.data.get("competitor_hospitals", []):
            p = tf.add_paragraph()
            p.text = f"[{comp['name']}] 강점: {comp['strength']}"
            sp = tf.add_paragraph(); sp.text = f"차이점: {comp['difference']}"; sp.level = 1

    def create_slide_09(self):
        """9장 홈페이지 구조 진단"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 9)
        slide.shapes.title.text = "홈페이지 AEO 기술 구조 진단"
        self._add_image_if_exists(slide, 9, left=1, top=4, width=4)
        
        diag = self.data.get("homepage_diagnosis", {})
        tf = slide.placeholders[1].text_frame
        tf.text = f"메인: {diag.get('main_page', '')}"
        p = tf.add_paragraph(); p.text = f"진료페이지: {diag.get('service_page', '')}"
        p = tf.add_paragraph(); p.text = f"Schema: {diag.get('schema_status', '')}"

    def create_slide_10(self):
        """10장 핵심 문제 3가지"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 10)
        slide.shapes.title.text = "AEO 도약을 위한 핵심 과제"
        tf = slide.placeholders[1].text_frame
        tf.text = ""
        for i, issue in enumerate(self.data.get("top_issues", [])[:3]):
            p = tf.add_paragraph()
            p.text = f"Issue {i+1}: {issue}"
            p.font.bold = True

    def create_slide_15(self):
        """15장 해결 전략 개요"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 15)
        slide.shapes.title.text = "Phase 2: AEO 해결 전략 (Proposal)"
        tf = slide.placeholders[1].text_frame
        tf.text = "전략 핵심 방향:"
        for plan in self.data.get("implementation_plan", {}).get("summary", []):
            p = tf.add_paragraph(); p.text = f"• {plan}"; p.level = 1

    def create_slide_16(self):
        """16장 핵심질문 1개 전략"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 16)
        slide.shapes.title.text = "The Core Question: 브랜드 검색 축 설정"
        core = self.data.get("core_question", {})
        tf = slide.placeholders[1].text_frame
        p = tf.add_paragraph()
        p.text = f"Q. \"{core.get('question', '')}\""
        p.font.size = Pt(28)
        p.font.bold = True
        
        p = tf.add_paragraph(); p.text = ""
        p = tf.add_paragraph(); p.text = "설정 근거:"
        for reason in core.get("why_this_question", []):
            sp = tf.add_paragraph(); sp.text = f"• {reason}"; sp.level = 1

    def create_slide_17(self):
        """17장 FAQ 예시 구조"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 17)
        slide.shapes.title.text = "AEO 인용을 위한 FAQ 구조화 예시"
        tf = slide.placeholders[1].text_frame
        tf.text = "카테고리별 대표 질문 설계:"
        
        for cat in self.data.get("faq_categories", []):
            p = tf.add_paragraph(); p.text = f"[{cat['category']}]"
            p.font.bold = True
            for q in cat.get("examples", [])[:1]: # 대표 1개만
                sp = tf.add_paragraph(); sp.text = f" - {q}"; sp.level = 1

    def create_slide_18(self):
        """18장 기술 구현 구조"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 18)
        slide.shapes.title.text = "홈페이지 기술 최적화 레이어"
        tf = slide.placeholders[1].text_frame
        tf.text = ""
        for step in self.data.get("implementation_plan", {}).get("tech_structure", []):
            p = tf.add_paragraph(); p.text = f"Step: {step}"

    def create_slide_19(self):
        """19장 실행 프로세스"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 19)
        slide.shapes.title.text = "AEO 구축 실행 타임라인"
        tf = slide.placeholders[1].text_frame
        tf.text = " -> ".join(self.data.get("implementation_plan", {}).get("process_steps", []))

    def create_slide_20(self):
        """20장 결론 및 제안"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._add_top_accent_bar(slide)
        self._add_footer_and_slide_number(slide, 20)
        slide.shapes.title.text = "결론: AI 검색 시대에 앞서가는 병원"
        tf = slide.placeholders[1].text_frame
        tf.text = "기대 효과:"
        for effect in self.data.get("expected_effects", []):
            p = tf.add_paragraph(); p.text = f"• {effect}"; p.level = 1
            
        p = tf.add_paragraph(); p.text = ""
        p = tf.add_paragraph()
        p.text = self.data.get("closing_message", "")
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
