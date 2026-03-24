import os
import json
from pptx import Presentation
from pptx.util import Inches

def build_ppt_v2():
    print("--- Start Building High-Fidelity PPT (v2) ---")
    
    # 1. 데이터 로드
    with open("0_input/hospital_data.json", "r", encoding="utf-8") as f:
        data = json.load(f)
    
    hospital_id = data.get("hospital_id", "unknown")
    hospital_name = data.get("hospital_name", "병원")
    output_dir = "5_output"
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    ppt_filename = f"{hospital_name}_AEO_Report_Premium.pptx"
    ppt_path = os.path.join(output_dir, ppt_filename)
    
    # 2. PPT 생성
    prs = Presentation()
    # 16:9 셋팅 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    render_dir = f"1_assets/hospitals/{hospital_id}/renders"
    
    # 슬라이드 순서 (1~20)
    for i in range(1, 21):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        
        # 렌더링된 이미지가 있는지 확인
        render_path = os.path.join(render_dir, f"slide_{i:02d}.png")
        
        if os.path.exists(render_path):
            print(f"Adding premium render for Slide {i}")
            # 전체 화면으로 이미지 삽입
            slide.shapes.add_picture(render_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            # 아직 렌더링되지 않은 슬라이드는 텍스트만 표시
            left = top = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, prs.slide_width - Inches(2), Inches(1))
            tf = txBox.text_frame
            tf.text = f"Slide {i}: [Template Needed]\n{hospital_name} AEO Strategy"
            print(f"Slide {i} is missing render, adding placeholder.")

    # 3. 저장
    prs.save(ppt_path)
    print(f"\n--- Success! Premium PPT Created: {ppt_path} ---")

if __name__ == "__main__":
    build_ppt_v2()
